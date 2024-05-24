import hashlib
import subprocess
import os

try:
    from pptx import Presentation
except ImportError:
    raise ImportError(
        'PPT file requires extra dependencies. Install with `pip install python-pptx`'
    ) from None
from embedchain.helpers.json_serializable import register_deserializable
from embedchain.loaders.base_loader import BaseLoader
from embedchain.utils.misc import clean_string

@register_deserializable
class PPTLoader(BaseLoader):
    def load_data(self, url):
        """Load data from a .pptx file."""
        output = []
        meta_data = {}
        extension = self.file_extension(url)
        tmp_file = url
        if extension == "ppt":
            self.convert_ppt_to_pptx(url)
            tmp_file = url+"x"
        content = clean_string(self.extract_text_from_ppt(tmp_file))

        meta_data["url"] = url
        output.append({"content": content, "meta_data": meta_data})
        doc_id = hashlib.sha256((content).encode()).hexdigest()
        os.remove(tmp_file)
        return {
            "doc_id": doc_id,
            "data": output
        }
    
    def extract_text_from_ppt(self, ppt_file_path):
        presentation = Presentation(ppt_file_path)

        all_text = []

        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_text = []

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            if slide_text:
                all_text.append(f"Slide {slide_number}:")
                all_text.extend(slide_text)
                all_text.append("\n")

        return "\n".join(all_text)
    
    def file_extension(self,filename):
        if '.' not in filename:
            return ''
        return filename.rsplit('.', 1)[1].lower()

    def convert_ppt_to_pptx(self,ppt_file):
        try:
            # 调用 unoconv 命令行工具进行转换
            subprocess.run(['unoconv', '-f', 'pptx', ppt_file], check=True)
            print("Conversion completed successfully.")
        except subprocess.CalledProcessError as e:
            print("Conversion failed:", e)